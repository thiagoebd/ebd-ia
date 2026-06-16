"""Gerador de PPTX com identidade visual EBD — v2.

Tipos: cover, intro, kpi_grid, stat_callout, table, bullets, quote_dark, closing.
Componentes consistentes: eyebrow + logo top-right + footer rule + cards border-top.
"""
from __future__ import annotations

import uuid
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

EBD_RED         = RGBColor(0xC8, 0x10, 0x2E)
WHITE           = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK       = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_MED        = RGBColor(0x33, 0x33, 0x33)
TEXT_GRAY       = RGBColor(0x66, 0x66, 0x66)
TEXT_LIGHT      = RGBColor(0x99, 0x99, 0x99)
DIVIDER         = RGBColor(0xE0, 0xE0, 0xE0)
SUBTLE_BG       = RGBColor(0xF9, 0xF9, 0xF9)
GROUP_BG        = RGBColor(0x1A, 0x1A, 0x1A)
MARK_OK         = RGBColor(0x2E, 0x7D, 0x32)
MARK_MID        = RGBColor(0xF3, 0x9C, 0x12)
MARK_WARN       = RGBColor(0xD9, 0x53, 0x4F)
RED_TINT        = RGBColor(0xFF, 0xCC, 0xCC)

FONT_FAMILY = "Calibri"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_solid_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def _add_text(slide, left, top, width, height, text, *,
              font_size=14, bold=False, italic=False,
              color=TEXT_DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = FONT_FAMILY; f.size = Pt(font_size)
    f.bold = bold; f.italic = italic
    f.color.rgb = color
    return box


def _add_runs(slide, left, top, width, height, runs, *,
              anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.line_spacing = line_spacing
    for r in runs:
        run = p.add_run()
        run.text = r["text"]
        f = run.font
        f.name = FONT_FAMILY
        f.size = Pt(r.get("size", 14))
        f.bold = r.get("bold", False)
        f.italic = r.get("italic", False)
        f.color.rgb = r.get("color", TEXT_DARK)
    return box


def _add_rect(slide, left, top, width, height, fill=WHITE, line=None, line_width=0.5):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _set_solid_fill(rect, fill)
    if line is not None:
        rect.line.color.rgb = line
        rect.line.width = Pt(line_width)
    return rect


def _eyebrow(slide, label):
    _add_text(slide, Inches(0.6), Inches(0.45), Inches(8.0), Inches(0.3),
              label.upper(), font_size=10, bold=True, color=TEXT_GRAY)
    _add_rect(slide, Inches(0.6), Inches(0.78), Inches(0.55), Pt(2.5), fill=EBD_RED)


def _logo_textual(slide):
    _add_text(slide, Inches(11.0), Inches(0.45), Inches(1.73), Inches(0.32),
              "EBD GRUPO",
              font_size=14, bold=True, color=EBD_RED, align=PP_ALIGN.RIGHT)
    _add_text(slide, Inches(11.0), Inches(0.78), Inches(1.73), Inches(0.25),
              "DESDE 1977",
              font_size=8, bold=True, color=TEXT_GRAY, align=PP_ALIGN.RIGHT)


def _footer_rule(slide, footer_left, page_num, total_pages):
    _add_rect(slide, Inches(0.6), Inches(7.05),
              Inches(12.13), Emu(8000), fill=DIVIDER)
    _add_text(slide, Inches(0.6), Inches(7.15), Inches(10), Inches(0.3),
              footer_left,
              font_size=9, color=TEXT_GRAY)
    page_text = f"{page_num:02d}" if total_pages is None else f"{page_num:02d} / {total_pages:02d}"
    _add_text(slide, Inches(11.0), Inches(7.15), Inches(1.73), Inches(0.3),
              page_text,
              font_size=10, bold=True, color=EBD_RED, align=PP_ALIGN.RIGHT)


def _content_chrome(slide, *, eyebrow, footer_left, page_num, total_pages):
    _eyebrow(slide, eyebrow)
    _logo_textual(slide)
    _footer_rule(slide, footer_left, page_num, total_pages)


def _card_white(slide, left, top, width, height):
    """Card branco com border-top vermelho 4px + borda fina cinza."""
    _add_rect(slide, left, top, width, height, fill=WHITE,
              line=DIVIDER, line_width=0.5)
    _add_rect(slide, left, top, width, Pt(4), fill=EBD_RED)


def _card_red(slide, left, top, width, height):
    """Card vermelho sólido."""
    _add_rect(slide, left, top, width, height, fill=EBD_RED)


# ─── SLIDES ──────────────────────────────────────────────────────────────

def slide_cover(prs, *, title, subtitle, footer, eyebrow_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bar_w = Inches(3.6)
    _add_rect(slide, SLIDE_W - bar_w, Emu(0), bar_w, SLIDE_H, fill=EBD_RED)

    _add_text(slide, Inches(9.85), Inches(0.5), Inches(3.2), Inches(0.35),
              "EBD GRUPO",
              font_size=14, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    _add_text(slide, Inches(9.85), Inches(0.82), Inches(3.2), Inches(0.25),
              "DESDE 1977",
              font_size=8, bold=True, color=RED_TINT, align=PP_ALIGN.RIGHT)

    if eyebrow_label:
        _add_text(slide, Inches(0.85), Inches(2.05), Inches(8.5), Inches(0.3),
                  eyebrow_label.upper(),
                  font_size=10, bold=True, color=TEXT_GRAY)

    _add_rect(slide, Inches(0.85), Inches(2.45), Inches(0.7), Pt(3), fill=EBD_RED)

    _add_text(slide, Inches(0.85), Inches(2.75), Inches(8.5), Inches(2.3),
              title,
              font_size=42, bold=True, color=TEXT_DARK, line_spacing=1.05)

    _add_text(slide, Inches(0.85), Inches(5.1), Inches(8.5), Inches(0.85),
              subtitle,
              font_size=20, bold=True, color=EBD_RED)

    _add_rect(slide, Inches(0.85), Inches(6.4), Inches(8.5), Emu(8000), fill=DIVIDER)

    _add_text(slide, Inches(0.85), Inches(6.55), Inches(8.5), Inches(0.4),
              footer,
              font_size=10, color=TEXT_GRAY)


def slide_intro(prs, *, eyebrow, title, lead, bullets,
                page_num, total_pages, footer_left):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _content_chrome(slide, eyebrow=eyebrow, footer_left=footer_left,
                    page_num=page_num, total_pages=total_pages)

    _add_text(slide, Inches(0.6), Inches(1.3), Inches(12.13), Inches(0.9),
              title,
              font_size=36, bold=True, color=TEXT_DARK, line_spacing=1.1)

    _add_text(slide, Inches(0.6), Inches(2.35), Inches(12.13), Inches(0.7),
              lead,
              font_size=15, color=TEXT_MED, line_spacing=1.4)

    top = Inches(3.3)
    item_h = Inches(0.6)
    gap = Inches(0.15)
    for i, b in enumerate(bullets):
        y = top + (item_h + gap) * i
        _add_text(slide, Inches(0.6), y, Inches(0.9), item_h,
                  f"{i+1:02d}",
                  font_size=28, bold=True, color=EBD_RED)
        _add_text(slide, Inches(1.7), y + Inches(0.05),
                  Inches(10.9), item_h,
                  b,
                  font_size=15, color=TEXT_DARK)


def slide_kpi_grid(prs, *, eyebrow, title, subtitle, kpis,
                   page_num, total_pages, footer_left):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _content_chrome(slide, eyebrow=eyebrow, footer_left=footer_left,
                    page_num=page_num, total_pages=total_pages)

    _add_text(slide, Inches(0.6), Inches(1.3), Inches(12.13), Inches(0.85),
              title,
              font_size=32, bold=True, color=TEXT_DARK)

    if subtitle:
        _add_text(slide, Inches(0.6), Inches(2.3), Inches(12.13), Inches(0.4),
                  subtitle,
                  font_size=13, color=TEXT_GRAY)

    n = len(kpis)
    if n == 0:
        return
    top_cards = Inches(3.0)
    card_h = Inches(3.5)
    total_w = Inches(12.13)
    gap = Inches(0.25)
    card_w = Emu(int((total_w - gap * (n - 1)) / n))

    for i, k in enumerate(kpis):
        x = Inches(0.6) + (card_w + gap) * i
        is_hi = k.get("highlighted", False)
        if is_hi:
            _card_red(slide, x, top_cards, card_w, card_h)
            label_c, value_c, desc_c = RED_TINT, WHITE, WHITE
        else:
            _card_white(slide, x, top_cards, card_w, card_h)
            label_c, value_c, desc_c = TEXT_GRAY, EBD_RED, TEXT_MED

        _add_text(slide, x + Inches(0.35), top_cards + Inches(0.4),
                  card_w - Inches(0.7), Inches(0.4),
                  k["label"].upper(),
                  font_size=11, bold=True, color=label_c)

        _add_text(slide, x + Inches(0.35), top_cards + Inches(0.95),
                  card_w - Inches(0.7), Inches(1.6),
                  k["value"],
                  font_size=36, bold=True, color=value_c)

        _add_text(slide, x + Inches(0.35), top_cards + Inches(2.5),
                  card_w - Inches(0.7), Inches(0.9),
                  k.get("description", ""),
                  font_size=11, color=desc_c, line_spacing=1.35)


def slide_stat_callout(prs, *, eyebrow, title, subtitle, stats,
                       page_num, total_pages, footer_left):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _content_chrome(slide, eyebrow=eyebrow, footer_left=footer_left,
                    page_num=page_num, total_pages=total_pages)

    _add_text(slide, Inches(0.6), Inches(1.3), Inches(12.13), Inches(0.85),
              title,
              font_size=32, bold=True, color=TEXT_DARK)

    if subtitle:
        _add_text(slide, Inches(0.6), Inches(2.3), Inches(12.13), Inches(0.4),
                  subtitle,
                  font_size=13, color=TEXT_GRAY)

    n = len(stats)
    if n == 0 or n > 3:
        return
    top_cards = Inches(3.0)
    card_h = Inches(3.5)
    total_w = Inches(12.13)
    gap = Inches(0.2)
    card_w = Emu(int((total_w - gap * (n - 1)) / n))

    for i, s in enumerate(stats):
        x = Inches(0.6) + (card_w + gap) * i
        is_hi = s.get("highlighted", False)
        if is_hi:
            _add_rect(slide, x, top_cards, card_w, card_h, fill=EBD_RED)
            label_c, value_c, desc_c = RED_TINT, WHITE, WHITE
        else:
            _add_rect(slide, x, top_cards, card_w, card_h, fill=SUBTLE_BG)
            label_c, value_c, desc_c = TEXT_GRAY, EBD_RED, TEXT_DARK

        _add_text(slide, x + Inches(0.4), top_cards + Inches(0.4),
                  card_w - Inches(0.8), Inches(0.4),
                  s["label"].upper(),
                  font_size=10, bold=True, color=label_c)

        # BIG NUMBER — 48pt cabe "R$ 4,82M" sem quebrar em card de ~4" largura
        _add_text(slide, x + Inches(0.4), top_cards + Inches(0.95),
                  card_w - Inches(0.8), Inches(1.7),
                  s["value"],
                  font_size=48, bold=True, color=value_c, line_spacing=1.0)

        _add_text(slide, x + Inches(0.4), top_cards + Inches(2.6),
                  card_w - Inches(0.8), Inches(0.85),
                  s.get("description", ""),
                  font_size=12, color=desc_c, line_spacing=1.3)


def slide_table(prs, *, eyebrow, title, subtitle, columns, rows,
                page_num, total_pages, footer_left):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _content_chrome(slide, eyebrow=eyebrow, footer_left=footer_left,
                    page_num=page_num, total_pages=total_pages)

    _add_text(slide, Inches(0.6), Inches(1.3), Inches(12.13), Inches(0.7),
              title,
              font_size=26, bold=True, color=TEXT_DARK)

    if subtitle:
        _add_text(slide, Inches(0.6), Inches(2.0), Inches(12.13), Inches(0.35),
                  subtitle,
                  font_size=11, color=TEXT_GRAY)
        table_top = Inches(2.45)
    else:
        table_top = Inches(2.05)

    table_w = Inches(12.13)
    widths_spec = [c.get("width", 1.0) for c in columns]
    total_spec = sum(widths_spec)
    col_widths = [Emu(int(table_w * (w / total_spec))) for w in widths_spec]

    row_h = Inches(0.35)
    header_h = Inches(0.44)
    group_h = Inches(0.3)

    x = Inches(0.6)
    _add_rect(slide, x, table_top, table_w, header_h, fill=EBD_RED)
    cur_x = x
    for col, w in zip(columns, col_widths):
        _add_text(slide, cur_x + Inches(0.18), table_top + Inches(0.08),
                  w - Inches(0.36), header_h - Inches(0.16),
                  col["label"],
                  font_size=11, bold=True, color=WHITE,
                  align=_align_for(col.get("type", "text")))
        cur_x += w

    cursor_y = table_top + header_h
    row_idx_visible = 0
    for row in rows:
        group_label = row.get("__group_start__")
        if group_label:
            _add_rect(slide, x, cursor_y, table_w, group_h, fill=GROUP_BG)
            _add_text(slide, x + Inches(0.25), cursor_y + Inches(0.06),
                      table_w - Inches(0.5), group_h - Inches(0.12),
                      group_label.upper(),
                      font_size=9, bold=True, color=WHITE)
            cursor_y += group_h
            row_idx_visible = 0
            continue

        if row_idx_visible % 2 == 1:
            _add_rect(slide, x, cursor_y, table_w, row_h, fill=SUBTLE_BG)
        cur_x = x
        for col, w in zip(columns, col_widths):
            raw = row.get(col["key"], "")
            mark = None
            if isinstance(raw, dict):
                mark = raw.get("mark")
                text = _format_value(raw.get("text", ""), col.get("type", "text"))
            else:
                text = _format_value(raw, col.get("type", "text"))

            is_first_col = (col is columns[0])
            color = EBD_RED if row.get("__highlight__") else TEXT_DARK

            text_x = cur_x + Inches(0.18)
            if mark in ("ok", "mid", "warn"):
                mark_colors = {"ok": MARK_OK, "mid": MARK_MID, "warn": MARK_WARN}
                _add_text(slide, cur_x + Inches(0.18), cursor_y + Inches(0.05),
                          Inches(0.3), row_h - Inches(0.12),
                          "●", font_size=14, bold=True, color=mark_colors[mark])
                text_x = cur_x + Inches(0.48)

            _add_text(slide, text_x, cursor_y + Inches(0.07),
                      w - (text_x - cur_x) - Inches(0.18), row_h - Inches(0.14),
                      text,
                      font_size=11, bold=is_first_col, color=color,
                      align=_align_for(col.get("type", "text")))
            cur_x += w

        cursor_y += row_h
        row_idx_visible += 1


def slide_bullets(prs, *, eyebrow, title, subtitle, items,
                  page_num, total_pages, footer_left):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _content_chrome(slide, eyebrow=eyebrow, footer_left=footer_left,
                    page_num=page_num, total_pages=total_pages)

    _add_text(slide, Inches(0.6), Inches(1.3), Inches(12.13), Inches(0.85),
              title,
              font_size=30, bold=True, color=TEXT_DARK)

    y_start = Inches(2.35)
    if subtitle:
        _add_text(slide, Inches(0.6), Inches(2.3), Inches(12.13), Inches(0.4),
                  subtitle,
                  font_size=13, color=TEXT_GRAY)
        y_start = Inches(2.85)

    item_h = Inches(0.65)
    gap = Inches(0.18)
    for i, item in enumerate(items):
        y = y_start + (item_h + gap) * i
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(0.7), y + Inches(0.22),
                                     Inches(0.16), Inches(0.16))
        _set_solid_fill(dot, EBD_RED)

        if isinstance(item, str):
            _add_text(slide, Inches(1.1), y, Inches(11.5), item_h, item,
                      font_size=15, color=TEXT_DARK)
        else:
            _add_runs(slide, Inches(1.1), y + Inches(0.05), Inches(11.5), item_h,
                      [
                          {"text": item.get("label", "") + ": ", "size": 15, "bold": True, "color": TEXT_DARK},
                          {"text": item.get("text", ""), "size": 15, "color": TEXT_MED},
                      ])


def slide_quote_dark(prs, *, eyebrow, quote, attribution,
                     page_num, total_pages, footer_left):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _content_chrome(slide, eyebrow=eyebrow, footer_left=footer_left,
                    page_num=page_num, total_pages=total_pages)

    block_top = Inches(2.5)
    block_h = Inches(3.5)
    _add_rect(slide, Inches(0.6), block_top, Inches(12.13), block_h, fill=TEXT_DARK)

    _add_text(slide, Inches(1.2), block_top + Inches(0.7),
              Inches(11.0), Inches(2.0),
              f'"{quote}"',
              font_size=30, color=WHITE, line_spacing=1.35)

    _add_text(slide, Inches(1.2), block_top + Inches(2.7),
              Inches(11.0), Inches(0.4),
              attribution.upper(),
              font_size=10, bold=True, color=RED_TINT)


def slide_closing(prs, *, quote, attribution, contact=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar_w = Inches(3.6)
    _add_rect(slide, SLIDE_W - bar_w, Emu(0), bar_w, SLIDE_H, fill=EBD_RED)

    _add_text(slide, Inches(9.85), Inches(0.5), Inches(3.2), Inches(0.35),
              "EBD GRUPO",
              font_size=14, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    _add_text(slide, Inches(9.85), Inches(0.82), Inches(3.2), Inches(0.25),
              "DESDE 1977",
              font_size=8, bold=True, color=RED_TINT, align=PP_ALIGN.RIGHT)

    _add_text(slide, Inches(0.85), Inches(2.0), Inches(1.5), Inches(1.0),
              '\u201C',
              font_size=80, color=DIVIDER, bold=True)

    _add_text(slide, Inches(0.85), Inches(3.0), Inches(8.5), Inches(1.6),
              quote,
              font_size=28, bold=True, italic=True, color=TEXT_DARK,
              line_spacing=1.2)

    _add_rect(slide, Inches(0.85), Inches(4.8), Inches(0.7), Pt(3), fill=EBD_RED)

    _add_text(slide, Inches(0.85), Inches(5.0), Inches(8.5), Inches(0.5),
              attribution,
              font_size=16, bold=True, color=EBD_RED)

    if contact:
        _add_text(slide, Inches(0.85), Inches(6.7), Inches(8.5), Inches(0.4),
                  contact,
                  font_size=10, color=TEXT_GRAY)


def _format_value(value, kind):
    if value is None or value == "":
        return ""
    if kind == "money":
        try:
            v = float(value)
            return f"R$ {v:_.2f}".replace(".", ",").replace("_", ".")
        except (TypeError, ValueError):
            return str(value)
    if kind == "int":
        try:
            return f"{int(value):_}".replace("_", ".")
        except (TypeError, ValueError):
            return str(value)
    if kind == "percent":
        try:
            return f"{float(value):.1f}%".replace(".", ",")
        except (TypeError, ValueError):
            return str(value)
    return str(value)


def _align_for(kind):
    if kind in ("money", "int", "percent"):
        return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT


def build_pptx(*, title, subtitle, slides,
               footer_author="EBD.ia",
               eyebrow_default="Relatório EBD.ia",
               output_path=None):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total_pages = len(slides)

    footer_left = f"Grupo EBD · {footer_author}"

    page = 1
    for s in slides:
        kind = s["kind"]
        eyebrow = s.get("eyebrow", eyebrow_default)

        if kind == "cover":
            slide_cover(prs, title=s.get("title", title),
                        subtitle=s.get("subtitle", subtitle),
                        footer=s.get("footer", footer_author),
                        eyebrow_label=s.get("eyebrow_label"))
        elif kind == "intro":
            page += 1
            slide_intro(prs, eyebrow=eyebrow,
                        title=s["title"], lead=s.get("lead", ""),
                        bullets=s.get("bullets", []),
                        page_num=page, total_pages=total_pages,
                        footer_left=footer_left)
        elif kind == "kpi_grid":
            page += 1
            slide_kpi_grid(prs, eyebrow=eyebrow,
                           title=s["title"], subtitle=s.get("subtitle"),
                           kpis=s["kpis"],
                           page_num=page, total_pages=total_pages,
                           footer_left=footer_left)
        elif kind == "stat_callout":
            page += 1
            slide_stat_callout(prs, eyebrow=eyebrow,
                               title=s["title"], subtitle=s.get("subtitle"),
                               stats=s["stats"],
                               page_num=page, total_pages=total_pages,
                               footer_left=footer_left)
        elif kind == "table":
            page += 1
            _cols = s["columns"]
            _rows = s["rows"]
            # Tolerância: columns como lista de strings -> dicts
            if _cols and isinstance(_cols[0], str):
                _cols = [{"key": str(c).lower().replace(" ", "_").replace("/", "_").replace("%", "pct").replace("ç","c").replace("ã","a").replace("á","a").replace("í","i").replace("ó","o").replace("ú","u").replace("é","e"),
                          "label": str(c), "type": "text"} for c in _cols]
            # Tolerância: rows como lista de listas -> dicts (pelos keys das colunas)
            if _rows and isinstance(_rows[0], (list, tuple)):
                _keys = [c["key"] for c in _cols]
                _rows = [dict(zip(_keys, r)) for r in _rows]
            slide_table(prs, eyebrow=eyebrow,
                        title=s["title"], subtitle=s.get("subtitle"),
                        columns=_cols, rows=_rows,
                        page_num=page, total_pages=total_pages,
                        footer_left=footer_left)
        elif kind == "bullets":
            page += 1
            _items = s.get("items") or s.get("bullets") or s.get("points") or s.get("texts") or []
            slide_bullets(prs, eyebrow=eyebrow,
                          title=s["title"], subtitle=s.get("subtitle"),
                          items=_items,
                          page_num=page, total_pages=total_pages,
                          footer_left=footer_left)
        elif kind == "quote_dark":
            page += 1
            slide_quote_dark(prs, eyebrow=eyebrow,
                             quote=s["quote"], attribution=s["attribution"],
                             page_num=page, total_pages=total_pages,
                             footer_left=footer_left)
        elif kind == "closing":
            slide_closing(prs, quote=s["quote"],
                          attribution=s["attribution"],
                          contact=s.get("contact"))
        else:
            raise ValueError(f"kind desconhecido: {kind}")

    if output_path is None:
        artifact_id = str(uuid.uuid4())
        output_path = Path(f"/tmp/{artifact_id}.pptx")
    else:
        artifact_id = output_path.stem

    prs.save(str(output_path))
    return artifact_id, output_path, output_path.name, output_path.stat().st_size


# ─── teste standalone Ferrero v2 ─────────────────────────────────────────

if __name__ == "__main__":
    SECTION = "Vendas · Ferrero · Junho 2026"

    test_slides = [
        {"kind": "cover",
         "title": "Vendas Ferrero — MTD jun/2026",
         "subtitle": "Faturamento Líquido por SKU · Visão Brasil",
         "eyebrow_label": "Relatório de Vendas · Diretoria Comercial",
         "footer": "Diretoria Comercial — Grupo EBD | EBD.ia | 16/06/2026"},

        {"kind": "intro",
         "eyebrow": f"01 — {SECTION}",
         "title": "O que vamos ver",
         "lead": "Análise consolidada do desempenho da linha Ferrero no mês corrente, "
                 "com foco em SKUs de maior giro e oportunidades.",
         "bullets": [
             "Resultados consolidados do mês — faturamento, pedidos, share no portfólio",
             "Composição do mix Ferrero ativo no mês vs portfólio total",
             "Top 10 SKUs Ferrero por faturamento líquido, agrupados por categoria",
             "Principais leituras e pontos de atenção pela equipe comercial",
         ]},

        {"kind": "stat_callout",
         "eyebrow": f"02 — {SECTION} · Visão geral",
         "title": "Resultados Consolidados",
         "subtitle": "Linha Ferrero · MTD jun/2026 · Visão Brasil",
         "stats": [
             {"label": "Faturamento Líquido MTD", "value": "R$ 4,82M",
              "description": "Acumulado mês corrente, +12,3% vs mai/2026"},
             {"label": "Pedidos no mês", "value": "3.847",
              "description": "Ticket médio R$ 1.253 por pedido"},
             {"label": "Share no Faturamento BR", "value": "8,7%",
              "description": "Ferrero entre os top 5 fornecedores do mês",
              "highlighted": True},
         ]},

        {"kind": "kpi_grid",
         "eyebrow": f"03 — {SECTION} · Composição",
         "title": "Composição do Mix Ferrero",
         "subtitle": "Mix de SKUs ativos no mês vs portfólio total · jun/2026",
         "kpis": [
             {"label": "SKUs ativos", "value": "47",
              "description": "de 62 SKUs do portfólio Ferrero EBD"},
             {"label": "Categorias", "value": "5 / 6",
              "description": "Falta apenas categoria 'Sazonais Páscoa'"},
             {"label": "Filiais c/ Ferrero", "value": "89%",
              "description": "24 de 27 filiais ativas com pelo menos 1 venda"},
             {"label": "SKUs sem giro", "value": "15",
              "description": "Risco de descontinuação · revisar com comercial",
              "highlighted": True},
         ]},

        {"kind": "table",
         "eyebrow": f"04 — {SECTION} · Top SKUs",
         "title": "Top 10 SKUs Ferrero",
         "subtitle": "Ordenado por faturamento líquido · MTD jun/2026 · agrupado por categoria",
         "columns": [
             {"key": "rank",  "label": "#",                   "type": "int",     "width": 0.45},
             {"key": "sku",   "label": "Produto",             "type": "text",    "width": 3.2},
             {"key": "ean",   "label": "EAN",                 "type": "text",    "width": 1.3},
             {"key": "qtd",   "label": "Caixas",              "type": "int",     "width": 0.85},
             {"key": "liq",   "label": "Faturamento Líquido", "type": "money",   "width": 1.5},
             {"key": "share", "label": "% Ferrero",           "type": "percent", "width": 0.9},
         ],
         "rows": [
             {"__group_start__": "A · CHOCOLATES PREMIUM (Nutella + Rocher + Raffaello)"},
             {"rank": 1, "sku": "Nutella 650g",                   "ean": "7898024390015", "qtd": 4820, "liq": 723415.80, "share": 15.0, "__highlight__": True},
             {"rank": 2, "sku": "Ferrero Rocher T16 (200g)",      "ean": "7898024390084", "qtd": 3210, "liq": 612893.40, "share": 12.7},
             {"rank": 4, "sku": "Nutella 350g",                   "ean": "7898024390022", "qtd": 5640, "liq": 451200.30, "share":  9.4},
             {"rank": 6, "sku": "Ferrero Rocher T8 (100g)",       "ean": "7898024390077", "qtd": 2415, "liq": 348521.60, "share":  7.2},
             {"rank":10, "sku": "Raffaello T15 (150g)",           "ean": "7898024390121", "qtd": 1120, "liq": 175218.50, "share":  3.6},

             {"__group_start__": "B · LINHA KINDER"},
             {"rank": 3, "sku": "Kinder Bueno Barra 43g (Cx 30)", "ean": "7898024391050", "qtd": 2890, "liq": 489120.15, "share": 10.1},
             {"rank": 5, "sku": "Kinder Ovo Surpresa 20g (Cx 36)","ean": "7898024391012", "qtd": 1830, "liq": 392874.90, "share":  8.2},

             {"__group_start__": "C · CONFEITARIA & MENTA"},
             {"rank": 7, "sku": "Tic Tac Menta 16g (Cx 24)",      "ean": "7898024399812", "qtd": 8200, "liq": 287400.00, "share":  6.0},
         ]},

        {"kind": "bullets",
         "eyebrow": f"05 — {SECTION} · Leituras",
         "title": "Principais Leituras",
         "subtitle": "Pontos de atenção e oportunidades para a equipe comercial",
         "items": [
             {"label": "Nutella 650g concentra 15%", "text": "única SKU passando dos R$ 700K, com presença em 89% das filiais ativas — risco de ruptura no fim do mês."},
             {"label": "Linha Kinder forte", "text": "Bueno + Ovo + Chocolate somam R$ 1,08M (22,4% do total) — performance acima da média histórica."},
             {"label": "Rocher T8 caiu 8,5% vs mai/2026", "text": "tendência inversa ao T16 — revisar estratégia de mix nos pequenos varejos."},
             {"label": "Tic Tac com volume alto e ticket baixo", "text": "8,2K caixas mas apenas R$ 287K — bom giro mas oportunidade de bundle."},
             {"label": "15 SKUs sem venda no mês", "text": "investigar com a área de comercial se há descontinuação ou problema de cadastro."},
         ]},
    ]

    out = Path("/home/claude/pptx_test/vendas-ferrero-v2.pptx")
    artifact_id, fp, fn, sz = build_pptx(
        title="Vendas Ferrero — MTD jun/2026",
        subtitle="Faturamento por SKU · Visão Brasil",
        slides=test_slides,
        footer_author="Diretoria Comercial · 16/06/2026",
        eyebrow_default="Relatório de Vendas",
        output_path=out,
    )
    print(f"OK PPTX v2 gerado")
    print(f"   id:       {artifact_id}")
    print(f"   path:     {fp}")
    print(f"   filename: {fn}")
    print(f"   size:     {sz} bytes ({sz/1024:.1f} KB)")
